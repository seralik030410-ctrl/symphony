"""Trusted offline parser entry point. Runs ONLY inside the restricted container."""
import json
from pathlib import Path
import sys
import zipfile
from contextlib import redirect_stdout

MAX_CHARS = 2_000_000


def extract(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".docx", ".pptx", ".xlsx"}:
        with zipfile.ZipFile(path) as archive:
            members = archive.infolist()
            if len(members) > 5000 or sum(item.file_size for item in members) > 80_000_000:
                raise ValueError("Office archive exceeds the 80 MB expanded/5000 entry limit")
            if any(item.file_size > 30_000_000 or item.flag_bits & 1 for item in members):
                raise ValueError("Encrypted or oversized Office archive member")
    parts = []
    total = 0

    def add(value):
        nonlocal total
        text = str(value)
        total += len(text) + 1
        if total > MAX_CHARS:
            raise ValueError("Extracted text exceeds 2 million characters; split the source into smaller files")
        parts.append(text)

    if suffix == ".pdf":
        import pymupdf
        with pymupdf.open(path) as document:
            if document.needs_pass or len(document) > 1000:
                raise ValueError("PDF is encrypted or exceeds 1000 pages")
            for index, page in enumerate(document):
                add(f"[Page {index+1}]\n{page.get_text('text')}")
    elif suffix == ".docx":
        from docx import Document
        document = Document(path)
        for paragraph in document.paragraphs:
            add(paragraph.text)
        for table in document.tables:
            for row in table.rows:
                add(" | ".join(cell.text for cell in row.cells))
    elif suffix == ".pptx":
        from pptx import Presentation
        presentation = Presentation(path)
        if len(presentation.slides) > 1000:
            raise ValueError("Presentation exceeds 1000 slides")
        for index, slide in enumerate(presentation.slides):
            add(f"[Slide {index+1}]")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    add(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        add(" | ".join(cell.text for cell in row.cells))
    elif suffix == ".xlsx":
        from openpyxl import load_workbook
        workbook = load_workbook(path, read_only=True, data_only=False, keep_links=False)
        try:
            if len(workbook.worksheets) > 64:
                raise ValueError("Workbook exceeds 64 sheets")
            for sheet in workbook.worksheets:
                if (sheet.max_row or 0) > 10000 or (sheet.max_column or 0) > 128:
                    raise ValueError("Sheet exceeds 10000 rows/128 columns; export the relevant range")
                add(f"[Sheet {sheet.title}]")
                for index, row in enumerate(sheet.iter_rows(values_only=True)):
                    if index >= 10000 or len(row) > 128:
                        raise ValueError("Sheet exceeds the extraction limit")
                    add("\t".join("" if value is None else str(value) for value in row))
        finally:
            workbook.close()
    else:
        raise ValueError("Unsupported extraction format")
    return "\n".join(parts)


if __name__ == "__main__":
    try:
        with redirect_stdout(sys.stderr):
            text = extract(Path(sys.argv[1]))
        print(json.dumps({"text": text}, ensure_ascii=False))
    except Exception as error:
        print(json.dumps({"error": str(error)[:1500]}))
        sys.exit(1)
