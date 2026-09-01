"""Trusted, deterministic renderer implementation. Input is data, never source code."""
from __future__ import annotations

import base64
import io
import json
import os
from pathlib import Path
import shutil
import subprocess
from xml.sax.saxutils import escape

from .formulas import calculate
from .schemas import ReportSpec, DocumentSpec, WorkbookSpec, SlideSpec

PALETTES = {"clean_report": "245E63", "finance_report": "294974", "technical_report": "34495E", "visual_brief": "9C4E35"}


def font_paths():
    for regular, bold in [
        ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        ("C:/Windows/Fonts/arial.ttf", "C:/Windows/Fonts/arialbd.ttf"),
        ("/Library/Fonts/Arial.ttf", "/Library/Fonts/Arial Bold.ttf"),
    ]:
        if Path(regular).is_file() and Path(bold).is_file(): return regular, bold
    raise ValueError("A Unicode font is required (DejaVu Sans or Arial)")


def image_bytes(spec):
    from PIL import Image
    raw = base64.b64decode(spec.png_base64, validate=True)
    with Image.open(io.BytesIO(raw)) as image:
        if image.width * image.height > 10_000_000 or max(image.size) > 6000:
            raise ValueError("PNG dimensions exceed 10 megapixels / 6000 pixels")
        image.verify()
    return raw


def chart_drawing(chart, accent="245E63"):
    from reportlab.graphics.shapes import Drawing, Rect, String, Line, PolyLine
    from reportlab.lib.colors import HexColor
    drawing = Drawing(450, 240)
    drawing.add(String(0, 222, chart.title[:70], fontName="SymphonyBold", fontSize=11, fillColor=HexColor("#18232B")))
    low, high = min(0, *chart.values), max(0, *chart.values)
    if high == low: high = low + 1
    y = lambda value: 45 + (value - low) / (high - low) * 145
    drawing.add(Line(25, y(0), 435, y(0), strokeColor=HexColor("#D0D8DD")))
    points = []
    step = 410 / len(chart.values)
    for index, (label, value) in enumerate(zip(chart.labels, chart.values)):
        x = 25 + step * (index + 0.5)
        if chart.kind == "bar":
            drawing.add(Rect(x - step * .3, min(y(0), y(value)), step * .6, max(.5, abs(y(value) - y(0))), fillColor=HexColor("#" + accent), strokeColor=None))
        points.extend([x, y(value)])
        drawing.add(String(x, 20, str(index + 1), textAnchor="middle", fontName="Symphony", fontSize=8))
        drawing.add(String(x, y(value) + (7 if value >= 0 else -12), f"{value:g}", textAnchor="middle", fontName="Symphony", fontSize=8))
    if chart.kind == "line":
        drawing.add(PolyLine(points, strokeColor=HexColor("#" + accent), strokeWidth=2))
    return drawing


def register_fonts():
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    regular, bold = font_paths()
    pdfmetrics.registerFont(TTFont("Symphony", regular))
    pdfmetrics.registerFont(TTFont("SymphonyBold", bold))
    pdfmetrics.registerFontFamily("Symphony", normal="Symphony", bold="SymphonyBold", italic="Symphony", boldItalic="SymphonyBold")


def chart_png(chart, accent):
    import pymupdf
    from reportlab.graphics import renderPDF
    register_fonts()
    raw = renderPDF.drawToString(chart_drawing(chart, accent))
    with pymupdf.open(stream=raw, filetype="pdf") as pdf:
        return pdf[0].get_pixmap(dpi=144).tobytes("png")


def render_pdf(spec: ReportSpec, target: Path):
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, KeepTogether
    register_fonts()
    accent = colors.HexColor("#" + PALETTES[spec.preset])
    body = ParagraphStyle("body", fontName="Symphony", fontSize=10, leading=15, textColor=colors.HexColor("#24343D"), spaceAfter=9, splitLongWords=True)
    heading = ParagraphStyle("heading", parent=body, fontName="SymphonyBold", fontSize=15, leading=20, spaceBefore=15, spaceAfter=10, keepWithNext=True, textColor=accent)
    title = ParagraphStyle("title", parent=heading, fontSize=27, leading=34, spaceBefore=0, spaceAfter=15)
    small = ParagraphStyle("small", parent=body, fontSize=8, leading=12)
    cell = ParagraphStyle("cell", parent=body, fontSize=8, leading=12, spaceAfter=0)
    def para(value, style=body): return Paragraph(escape(str(value if value is not None else "")).replace("\n", "<br/>"), style)
    story = [para(spec.title, title)]
    if spec.subtitle: story.append(para(spec.subtitle))
    for section in spec.sections:
        story.append(para(section.heading, heading))
        story.extend(para(text) for text in section.paragraphs)
        story.extend(para("• " + text) for text in section.bullets)
        if section.callout:
            block = Table([[para(section.callout)]], colWidths=[A4[0] - 104])
            block.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F0F5F5")), ("BOX", (0, 0), (-1, -1), .5, accent), ("LEFTPADDING", (0, 0), (-1, -1), 12), ("RIGHTPADDING", (0, 0), (-1, -1), 12), ("TOPPADDING", (0, 0), (-1, -1), 12)]))
            story.extend([block, Spacer(1, 10)])
        if section.table:
            table = section.table
            rows = [[para(v, cell) for v in table.columns]] + [[para(v, cell) for v in row] for row in table.rows]
            grid = Table(rows, colWidths=[(A4[0] - 104) / len(table.columns)] * len(table.columns), repeatRows=1, hAlign="LEFT")
            grid.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#DDEBEB")), ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F6F8F9")]), ("VALIGN", (0, 0), (-1, -1), "TOP"), ("TOPPADDING", (0, 0), (-1, -1), 7), ("BOTTOMPADDING", (0, 0), (-1, -1), 7), ("LINEBELOW", (0, 0), (-1, 0), .8, accent)]))
            story.extend([grid, Spacer(1, 12)])
        if section.chart:
            story.append(chart_drawing(section.chart, PALETTES[spec.preset]))
            story.append(para(" · ".join(f"{i + 1}. {label}" for i, label in enumerate(section.chart.labels)), small))
        if section.image:
            image = Image(io.BytesIO(image_bytes(section.image)))
            scale = min(480 / image.imageWidth, 320 / image.imageHeight, 1)
            image.drawWidth, image.drawHeight = image.imageWidth * scale, image.imageHeight * scale
            story.append(KeepTogether([image, para(section.image.caption, small)]))
    if spec.citations:
        story.append(para("Источники", heading))
        story.extend(para(f"{i + 1}. {text}", small) for i, text in enumerate(spec.citations))

    def footer(canvas, doc):
        if doc.page > 80: raise ValueError("Document exceeds 80 pages")
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor("#D8E1E4")); canvas.line(48, 39, A4[0] - 48, 39)
        canvas.setFont("Symphony", 8); canvas.setFillColor(colors.HexColor("#60727C"))
        canvas.drawString(48, 25, "Symphony · " + spec.preset.replace("_", " "))
        canvas.drawRightString(A4[0] - 48, 25, str(doc.page)); canvas.restoreState()
    SimpleDocTemplate(str(target), pagesize=A4, rightMargin=52, leftMargin=52, topMargin=48, bottomMargin=55, title=spec.title, author="Symphony").build(story, onFirstPage=footer, onLaterPages=footer)


def render_workbook(spec: WorkbookSpec, target: Path):
    from datetime import date
    from openpyxl import Workbook, load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter
    calculation = calculate(spec)
    workbook = Workbook(); workbook.remove(workbook.active)
    workbook.properties.title = spec.title
    formats = {"general": "General", "integer": "#,##0", "decimal": "#,##0.00", "currency": '#,##0.00;[Red](#,##0.00)', "percent": "0.0%", "date": "yyyy-mm-dd"}
    tables = []
    for sheet in spec.sheets:
        ws = workbook.create_sheet(sheet.name)
        ws.append([column.name for column in sheet.columns])
        for row_index, row in enumerate(sheet.rows, 2):
            for col_index, (column, value) in enumerate(zip(sheet.columns, row), 1):
                cell = ws.cell(row_index, col_index)
                cell.value = date.fromisoformat(value) if column.type == "date" and value else value
                if isinstance(value, str) and column.type != "date": cell.data_type = "s"
                cell.number_format = formats[column.format]
                cell.alignment = Alignment(vertical="top", wrap_text=True)
        for address, formula in sheet.formulas.items(): ws[address] = formula
        for cell in ws[1]:
            cell.data_type = "s"
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="245E63")
            cell.alignment = Alignment(wrap_text=True)
        for index, column in enumerate(sheet.columns, 1): ws.column_dimensions[get_column_letter(index)].width = column.width
        ws.freeze_panes = "A2"; ws.auto_filter.ref = ws.dimensions
        ws.sheet_view.showGridLines = False
        rows = []
        for index, row in enumerate(sheet.rows[:200], 2):
            rows.append([calculation["values"].get(f"{sheet.name}!{get_column_letter(c)}{index}", value) for c, value in enumerate(row, 1)])
        tables.append({"name": sheet.name, "columns": [c.model_dump() for c in sheet.columns], "rows": rows, "total_rows": len(sheet.rows), "formulas": sheet.formulas, "truncated": len(sheet.rows) > 200})
    workbook.save(target)
    reopened = load_workbook(target, data_only=False, keep_links=False)
    for sheet in spec.sheets:
        for address, formula in sheet.formulas.items():
            if reopened[sheet.name][address].value != formula: raise ValueError("Formula round-trip validation failed")
        for row in reopened[sheet.name]:
            for cell in row:
                if cell.data_type == "f" and cell.coordinate not in sheet.formulas: raise ValueError("Unexpected formula injection")
    reopened.close()
    return {"tables": tables, "calculation": calculation, "warnings": ["Excel recalculates formulas when opened; preview values come from Symphony's bounded calculator."] if calculation["formula_count"] else []}


def render_docx(spec: DocumentSpec, target: Path):
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    document = Document()
    normal = document.styles["Normal"]
    normal.font.name = "DejaVu Sans"; normal.font.size = Pt(10)
    normal.paragraph_format.space_after = Pt(8)
    for name in ["Title", "Heading 1", "Heading 2"]:
        document.styles[name].font.name = "DejaVu Sans"
        document.styles[name].font.color.rgb = RGBColor.from_string(PALETTES[spec.preset])
    for section in document.sections:
        section.header.paragraphs[0].text = spec.header
        section.footer.paragraphs[0].text = spec.footer
        section.top_margin = section.bottom_margin = Inches(.7)
    document.add_heading(spec.title, 0)
    if spec.subtitle: document.add_paragraph(spec.subtitle, "Subtitle")
    for section in spec.sections:
        document.add_heading(section.heading, 1)
        for text in section.paragraphs: document.add_paragraph(text)
        for text in section.bullets: document.add_paragraph(text, "List Bullet")
        if section.callout: document.add_paragraph(section.callout, "Quote")
        if section.table:
            table = document.add_table(rows=1, cols=len(section.table.columns)); table.style = "Light Shading Accent 1"
            for cell, text in zip(table.rows[0].cells, section.table.columns): cell.text = text
            for row in section.table.rows:
                for cell, value in zip(table.add_row().cells, row): cell.text = str(value if value is not None else "")
        if section.chart:
            document.add_picture(io.BytesIO(chart_png(section.chart, PALETTES[spec.preset])), width=Inches(6))
            document.add_paragraph(" · ".join(f"{i + 1}. {label}" for i, label in enumerate(section.chart.labels)), "Caption")
        if section.image:
            document.add_picture(io.BytesIO(image_bytes(section.image)), width=Inches(5))
            document.add_paragraph(section.image.caption, "Caption")
    if spec.citations:
        document.add_heading("Источники", 1)
        for i, text in enumerate(spec.citations, 1): document.add_paragraph(f"{i}. {text}")
    document.core_properties.title = spec.title
    document.save(target)
    Document(target)  # Reopen the actual OOXML, not a mock preview.


def render_slides(spec: SlideSpec, target: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    prs.core_properties.title = spec.title
    accent = PALETTES[spec.preset]
    def text(slide, value, x, y, w, h, size=22, bold=False, color="203039"):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = value; p.font.name = "DejaVu Sans"; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = RGBColor.from_string(color)
        return box
    for index, item in enumerate(spec.slides, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.1)); band.fill.solid(); band.fill.fore_color.rgb = RGBColor.from_string(accent); band.line.fill.background()
        text(slide, item.title, .65, .42, 12, .95, 30, True, accent)
        if item.subtitle: text(slide, item.subtitle, .65, 1.35, 12, .65, 15)
        if item.bullets:
            for i, bullet in enumerate(item.bullets): text(slide, "• " + bullet, .8, 2.1 + i * .7, 11.7, .65, 18)
        elif item.table:
            rows, cols = len(item.table.rows) + 1, len(item.table.columns)
            table = slide.shapes.add_table(rows, cols, Inches(.7), Inches(2.1), Inches(11.9), Inches(4.4)).table
            for r, values in enumerate([item.table.columns, *item.table.rows]):
                for c, value in enumerate(values):
                    cell = table.cell(r, c); cell.text = str(value if value is not None else "")
                    for p in cell.text_frame.paragraphs: p.font.size = Pt(12 if rows > 7 else 15); p.font.name = "DejaVu Sans"
        elif item.chart:
            data = CategoryChartData(); data.categories = item.chart.labels; data.add_series(item.chart.title, item.chart.values)
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED if item.chart.kind == "bar" else XL_CHART_TYPE.LINE, Inches(.7), Inches(2.1), Inches(11.9), Inches(4.45), data)
        elif item.image:
            raw = image_bytes(item.image)
            from PIL import Image
            with Image.open(io.BytesIO(raw)) as image: w, h = image.size
            scale = min(11.5 / w, 4.2 / h)
            slide.shapes.add_picture(io.BytesIO(raw), Inches((13.333 - w * scale) / 2), Inches(2.1), width=Inches(w * scale), height=Inches(h * scale))
            text(slide, item.image.caption, .7, 6.4, 11.9, .35, 10)
        text(slide, f"Symphony    /    {index:02d}", .7, 7.03, 12, .3, 10, color="677780")
        if item.notes: slide.notes_slide.notes_text_frame.text = item.notes
    prs.save(target)
    check = Presentation(target)
    if len(check.slides) != len(spec.slides): raise ValueError("Slide count validation failed")
    for slide in check.slides:
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                raise ValueError("Slide shape extends beyond canvas")


def office_pdf(target: Path) -> Path:
    binary = shutil.which("libreoffice") or shutil.which("soffice")
    if not binary: raise ValueError("LibreOffice is required for verified DOCX/PPTX preview; rebuild the document runtime")
    folder = target.parent / "office-preview"; folder.mkdir(exist_ok=True)
    result = subprocess.run([binary, "-env:UserInstallation=file:///tmp/symphony-office", "--headless", "--convert-to", "pdf", "--outdir", str(folder), str(target)], capture_output=True, timeout=90, check=False)
    pdf = folder / (target.stem + ".pdf")
    if result.returncode or not pdf.is_file(): raise ValueError("Office-to-PDF conversion failed: " + result.stderr.decode(errors="replace")[:500])
    final = target.parent / "preview.pdf"; shutil.move(pdf, final); folder.rmdir()
    return final


def inspect_pdf(target: Path, out: Path, expected_pages: int | None = None) -> dict:
    import pymupdf
    pages = []; problems = []
    with pymupdf.open(target) as document:
        if not 1 <= len(document) <= 80: raise ValueError("PDF must contain 1–80 pages")
        if expected_pages is not None and len(document) != expected_pages: raise ValueError("Converted slide/page count mismatch")
        for index, page in enumerate(document):
            bounds = page.rect + (-1, -1, 1, 1)
            for block in page.get_text("dict")["blocks"]:
                if block["type"] == 0 and not bounds.contains(pymupdf.Rect(block["bbox"])):
                    problems.append(f"Text outside page bounds: page {index + 1}")
            filename = f"page-{index + 1:03d}.png"
            page.get_pixmap(dpi=110, alpha=False).save(out / filename)
            pages.append({"file": filename, "width": round(page.rect.width), "height": round(page.rect.height), "characters": len(page.get_text())})
    if problems: raise ValueError("; ".join(problems[:5]))
    return {"pages": pages, "geometry": {"checked": True, "out_of_bounds_text": 0}, "warnings": ["Automatic geometry checks do not replace visual review of content, overlaps or font shaping."]}


def render_job(job: Path) -> dict:
    from .schemas import parse_spec
    source = json.loads((job / "source.json").read_text(encoding="utf-8"))
    recipe = json.loads((job / "recipe.json").read_text(encoding="utf-8"))
    format = recipe["format"]; spec = parse_spec(format, source)
    target = job / f"document.{format}"
    if format == "pdf": render_pdf(spec, target)
    elif format == "xlsx": result = render_workbook(spec, target)
    elif format == "docx": render_docx(spec, target)
    elif format == "pptx": render_slides(spec, target)
    if format != "xlsx":
        pdf = target if format == "pdf" else office_pdf(target)
        result = inspect_pdf(pdf, job, len(spec.slides) if format == "pptx" else None)
        if format != "pdf": result["preview_pdf"] = "preview.pdf"
    result.update({"valid": True, "schema": type(spec).__name__, "renderer": "symphony-documents-v1", "output": target.name, "title": spec.title, "format": format})
    from importlib.metadata import version
    result["dependencies"] = {name: version(name) for name in ["reportlab", "openpyxl", "python-docx", "python-pptx", "pymupdf", "pydantic"]}
    (job / "validation.json").write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    return result
