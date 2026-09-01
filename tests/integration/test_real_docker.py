"""Opt-in real-container tests; never silently substitute a fake runtime.

PowerShell: $env:SYMPHONY_RUN_DOCKER_TESTS='1'; .venv/Scripts/python -m pytest tests/integration/test_real_docker.py
"""
import asyncio
import os
import uuid

import httpx
import pytest

from backend.main import create_app
from backend.models.base import ModelStreamEvent, ToolCall
from backend.models.gateway import ModelGateway
from backend.sandbox.runtime import DockerSandboxRuntime
from backend.skills.store import SkillStore
from backend.storage.database import Database
from backend.tools.contracts import ToolContext, ToolError
from backend.tools.skills import RunSkillScriptInput, RunSkillScriptTool
from backend.tools.workspace import WorkspaceManager
from conftest import FakeAdapter, wait_for_final

pytestmark = pytest.mark.skipif(os.getenv("SYMPHONY_RUN_DOCKER_TESTS") != "1", reason="real Docker opt-in")


@pytest.mark.parametrize("suffix", [".pdf", ".docx", ".pptx", ".xlsx"])
async def test_real_isolated_document_extraction(tmp_path, suffix):
    from backend.agent.extraction import IsolatedExtractor
    marker = "CALIBRATION_73421"
    path = tmp_path / ("source" + suffix)
    if suffix == ".pdf":
        import pymupdf
        with pymupdf.open() as document:
            document.new_page().insert_text((50, 50), marker)
            document.save(path)
    elif suffix == ".docx":
        from docx import Document
        document = Document(); document.add_paragraph(marker); document.save(path)
    elif suffix == ".pptx":
        from pptx import Presentation
        presentation = Presentation(); presentation.slides.add_slide(presentation.slide_layouts[0]).shapes.title.text = marker; presentation.save(path)
    else:
        from openpyxl import Workbook
        workbook = Workbook(); workbook.active.append([marker, 42]); workbook.save(path)
    text = await IsolatedExtractor(make_runtime(tmp_path)).extract(path.read_bytes(), suffix)
    assert marker in text


async def test_real_parser_cancellation_removes_container(tmp_path, monkeypatch):
    from backend.agent.extraction import IsolatedExtractor
    runtime = make_runtime(tmp_path)
    extractor = IsolatedExtractor(runtime)
    original = extractor.arguments
    names = []
    def slow_arguments(source, name):
        names.append(name)
        args = original(source, name)
        return args[:-3] + ["python", "-c", "import time; time.sleep(120)"]
    monkeypatch.setattr(extractor, "arguments", slow_arguments)
    task = asyncio.create_task(extractor.extract(b"test", ".pdf"))
    for _ in range(30):
        await asyncio.sleep(0.1)
        if names:
            code, _ = await runtime._docker_control("inspect", names[0])
            if not code:
                break
    task.cancel()
    with pytest.raises(asyncio.CancelledError):
        await task
    assert names
    code, _ = await runtime._docker_control("inspect", names[0])
    assert code != 0


def make_runtime(tmp_path, **kwargs):
    return DockerSandboxRuntime(WorkspaceManager(tmp_path / "workspaces"), image="symphony-sandbox:stage3", **kwargs)


async def execute(runtime, command, *, session_id="a" * 32, timeout=30, on_output=None):
    return await runtime.execute(session_id=session_id, turn_id=uuid.uuid4().hex, command=command,
                                 cwd=".", timeout_seconds=timeout, network=False, on_output=on_output)


async def test_real_runtime_versions_and_isolation(tmp_path):
    runtime = make_runtime(tmp_path)
    source = '''import os, sys, socket, pathlib, pytest, openpyxl, docx, pptx, pypdf
assert sys.version_info[:2] == (3, 12)
assert os.getuid() == 10001
assert not pathlib.Path('/var/run/docker.sock').exists()
assert 'OPENAI_API_KEY' not in os.environ
try:
    pathlib.Path('/opt/forbidden').write_text('no')
except OSError:
    pass
else:
    raise AssertionError('root filesystem was writable')
try:
    socket.create_connection(('1.1.1.1', 443), timeout=1)
except OSError:
    pass
else:
    raise AssertionError('network was enabled')
print(sys.version, 'isolation passed')
'''
    runtime.workspaces.resolve("a" * 32, "check.py").write_text(source)
    result = await execute(runtime, "python check.py && node --version && npm --version && rg --version && curl --version && pdftotext -v")
    assert "isolation passed" in result.stdout
    assert "v22." in result.stdout


async def test_real_output_is_streamed_and_bounded(tmp_path):
    runtime = make_runtime(tmp_path, output_limit=1024)
    chunks = []
    async def capture(chunk):
        chunks.append(chunk)
    result = await execute(runtime, "python -c \"print('x'*20000)\"", on_output=capture)
    assert result.output_truncated
    assert len(result.stdout) == 1024
    assert sum(len(c["delta"]) for c in chunks) == 1024


async def test_real_skill_script_is_offline_readonly_and_workspace_scoped(tmp_path):
    runtime = make_runtime(tmp_path)
    database = Database(tmp_path / "skills.db")
    database.initialize()
    source = tmp_path / "source-skill"
    (source / "scripts").mkdir(parents=True)
    (source / "SKILL.md").write_text("---\nname: Runtime check\ndescription: Verify a local project offline\n---\n# Runtime check\n", encoding="utf-8")
    (source / "scripts" / "check.py").write_text(
        "from pathlib import Path\n"
        "Path('skill-output.txt').write_text('verified')\n"
        "try:\n Path('/skill/SKILL.md').write_text('mutated')\nexcept OSError:\n print('SKILL_READONLY')\nelse:\n raise SystemExit('skill mount writable')\n",
        encoding="utf-8",
    )
    skills = SkillStore(database, tmp_path / "managed-skills")
    installed = skills.install_folder(str(source), mode="explicit")
    tool = RunSkillScriptTool(skills, runtime)
    context = ToolContext(session_id="b" * 32, turn_id=uuid.uuid4().hex, selected_skill_ids={installed["id"]})
    result = await tool.execute(context, RunSkillScriptInput(skill_id=installed["id"], path="scripts/check.py"))
    assert "SKILL_READONLY" in result.output["stdout"]
    assert result.changed_files == ["skill-output.txt"]
    assert runtime.workspaces.resolve("b" * 32, "skill-output.txt").read_text() == "verified"
    assert skills.get(installed["id"])["skill_md"].startswith("---")


@pytest.mark.parametrize("cancel", [False, True])
async def test_real_timeout_and_stop_kill_child_tree(tmp_path, cancel):
    runtime = make_runtime(tmp_path)
    root = runtime.workspaces.session_root("a" * 32)
    (root / "child.py").write_text("import time\nfrom pathlib import Path\ntime.sleep(5)\nPath('escaped.txt').write_text('bad')\ntime.sleep(60)\n")
    (root / "parent.py").write_text("import subprocess,time\nfrom pathlib import Path\nsubprocess.Popen(['python','child.py'])\nPath('ready.txt').write_text('yes')\ntime.sleep(60)\n")
    task = asyncio.create_task(execute(runtime, "python parent.py", timeout=30 if cancel else 2))
    if cancel:
        for _ in range(100):
            if (root / "ready.txt").exists():
                break
            await asyncio.sleep(.1)
        assert (root / "ready.txt").exists()
        task.cancel()
        with pytest.raises(asyncio.CancelledError):
            await task
    else:
        with pytest.raises(ToolError) as error:
            await task
        assert error.value.code == "sandbox_timeout"
    code, containers = await runtime._docker_control("ps", "-aq", "--filter", f"label=com.symphony.instance={runtime.instance_id}")
    assert code == 0 and not containers.strip()
    await asyncio.sleep(5.1)
    assert not (root / "escaped.txt").exists()


async def test_real_restart_cleans_only_own_orphan_and_keeps_project(tmp_path):
    runtime = make_runtime(tmp_path)
    root = runtime.workspaces.session_root("a" * 32)
    (root / "source.txt").write_text("persistent")
    args = runtime.docker_arguments(workspace=root, container_name=f"symphony-qa-{uuid.uuid4().hex[:8]}",
                                    container_cwd="/workspace", command="sleep 60", network=False)
    code, orphan = await runtime._docker_control(*[args[0], "--detach", *args[1:]])
    assert code == 0
    restarted = make_runtime(tmp_path)
    await restarted.recover_orphans()
    code, _ = await restarted._docker_control("inspect", orphan.strip())
    assert code != 0
    assert (root / "source.txt").read_text() == "persistent"


class RealBuildAdapter(FakeAdapter):
    async def stream_chat(self, request):
        self.requests.append(request)
        step = len(self.requests)
        if step == 1:
            yield ModelStreamEvent(type="tool_call", tool_call=ToolCall("pkg", "fs.write", {"path": "package.json", "content": '{"scripts":{"build":"node build.cjs"}}'}))
            yield ModelStreamEvent(type="tool_call", tool_call=ToolCall("source", "fs.write", {"path": "build.cjs", "content": "const fs=require('node:fs');fs.mkdirSync('dist',{recursive:true});fs.writeFileSync('dist/index.html','<!doctype html><h1>Real Docker build</h1>');console.log('BUILD_OK');"}))
        elif step == 2:
            yield ModelStreamEvent(type="tool_call", tool_call=ToolCall("build", "sandbox.shell", {"command": "npm run build"}))
        elif step == 3:
            yield ModelStreamEvent(type="tool_call", tool_call=ToolCall("preview", "sandbox.preview", {"entry": "dist/index.html"}))
        else:
            yield ModelStreamEvent(type="text_delta", delta="Built and previewed.")


async def test_real_build_api_reload_and_restart(settings):
    adapter = RealBuildAdapter("ollama")
    app = create_app(settings, ModelGateway({"ollama": adapter}))
    async with app.router.lifespan_context(app):
        async with httpx.AsyncClient(transport=httpx.ASGITransport(app=app), base_url="http://test") as client:
            session = (await client.post("/api/sessions", json={})).json()
            created = (await client.post(f"/api/sessions/{session['id']}/turns", json={"content": "build site"})).json()
            turn = await wait_for_final(client, created["turn"]["id"], timeout=45)
            assert turn["status"] == "completed", turn
            events = (await client.get(f"/api/turns/{turn['id']}/events")).json()
            assert any(e["type"] == "tool.output_delta" and "BUILD_OK" in e["payload"]["delta"] for e in events)
            url = next(e["payload"]["preview_url"] for e in events if e["type"] == "preview.ready")
            assert "Real Docker build" in (await client.get(url)).text
    restarted = create_app(settings, ModelGateway({"ollama": adapter}))
    async with restarted.router.lifespan_context(restarted):
        async with httpx.AsyncClient(transport=httpx.ASGITransport(app=restarted), base_url="http://test") as client:
            saved = (await client.get(f"/api/sessions/{session['id']}")).json()
            assert saved["turns"][0]["status"] == "completed"
            assert "Real Docker build" in (await client.get(url)).text
            assert len((await client.get(f"/api/sessions/{session['id']}/snapshots")).json()) == 3
